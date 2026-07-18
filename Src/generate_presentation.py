from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_presentation(output_path):
    prs = Presentation()

    # Define common slide layouts
    title_slide_layout = prs.slide_layouts[0]
    bullet_slide_layout = prs.slide_layouts[1]
    
    # ----------------------------------------------------------------------
    # Slide 1: Title Slide
    # ----------------------------------------------------------------------
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Product Cannibalization & Demand Shift Analysis"
    subtitle.text = "Executive Summary & Business Insights\nPrepared by: Harish Chavan"

    # ----------------------------------------------------------------------
    # Slide 2: Executive Summary
    # ----------------------------------------------------------------------
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "Executive Summary"
    
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.text = "We analyzed 60,000 transaction records (Jan 2023 - Dec 2024) to evaluate the impact of launching Products 4 and 5 on our existing portfolio."
    
    p = tf.add_paragraph()
    p.text = "Key Finding: The launch of Products 4 and 5 generated ₹88.3M in new revenue, but caused a 19% sales drop in our legacy flagship, Product 6."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Net Impact: True portfolio growth is lower than raw revenue suggests due to internal demand shifting."
    p.level = 1

    # ----------------------------------------------------------------------
    # Slide 3: Cannibalization Metrics
    # ----------------------------------------------------------------------
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "Cannibalization Metrics"
    
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.text = "Product 6 experienced severe cannibalization after the new launches:"
    
    p = tf.add_paragraph()
    p.text = "Sales Volume Drop: Decreased by 19.3% in the post-launch period."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Revenue Loss: Over ₹25M in revenue shifted away from Product 6."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Customer Migration: 14.9% of loyal Product 6 buyers switched exclusively to Product 4."
    p.level = 1

    # ----------------------------------------------------------------------
    # Slide 4: Business Recommendations
    # ----------------------------------------------------------------------
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "Strategic Recommendations"
    
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.text = "To protect legacy margins while growing market share:"
    
    p = tf.add_paragraph()
    p.text = "1. Geographic Repositioning: Push Product 4 in the West (where it performs best) and protect Product 6 in the North."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "2. Price Differentiation: Slightly increase the price of Product 4 to widen the gap with Product 6, reducing direct overlap."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "3. Targeted Promotions: Offer loyalty discounts for Product 6 to retain the 14.9% of migrating customers."
    p.level = 1

    # Save
    prs.save(output_path)
    print(f"Presentation saved successfully to: {output_path}")

if __name__ == "__main__":
    create_presentation("Reports/Executive_Presentation.pptx")
